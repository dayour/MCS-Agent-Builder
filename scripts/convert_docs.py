#!/usr/bin/env python3
"""
Document Converter - Convert PDF, DOCX, PPTX to Markdown

Usage:
    python convert_docs.py <input_file> [output_file]
    python convert_docs.py --batch <input_folder> [output_folder]

Examples:
    python convert_docs.py document.pdf
    python convert_docs.py document.docx output.md
    python convert_docs.py --batch ./documents ./markdown
"""

import argparse
import sys
from pathlib import Path
from typing import Optional


def convert_docx_to_md(input_path: Path) -> str:
    """Convert DOCX to Markdown."""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("Install python-docx: pip install python-docx")

    doc = Document(input_path)
    lines = []

    for para in doc.paragraphs:
        style = para.style.name.lower() if para.style else ""
        text = para.text.strip()

        if not text:
            lines.append("")
            continue

        # Convert headings
        if "heading 1" in style:
            lines.append(f"# {text}")
        elif "heading 2" in style:
            lines.append(f"## {text}")
        elif "heading 3" in style:
            lines.append(f"### {text}")
        elif "heading 4" in style:
            lines.append(f"#### {text}")
        elif "list" in style:
            lines.append(f"- {text}")
        else:
            lines.append(text)

    # Handle tables
    for table in doc.tables:
        lines.append("")
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
            if i == 0:
                lines.append("|" + "|".join(["---"] * len(cells)) + "|")
        lines.append("")

    return "\n".join(lines)


def convert_pptx_to_md(input_path: Path) -> str:
    """Convert PPTX to Markdown."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("Install python-pptx: pip install python-pptx")

    prs = Presentation(input_path)
    lines = []

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        lines.append("")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                # Check if it's a title
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    if shape.placeholder_format.type == 1:  # Title
                        lines.append(f"### {text}")
                    else:
                        lines.append(text)
                else:
                    lines.append(text)

            # Handle tables in slides
            if shape.has_table:
                table = shape.table
                lines.append("")
                for j, row in enumerate(table.rows):
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
                    if j == 0:
                        lines.append("|" + "|".join(["---"] * len(cells)) + "|")
                lines.append("")

        lines.append("")
        lines.append("---")
        lines.append("")

    return "\n".join(lines)


def convert_pdf_to_md(input_path: Path) -> str:
    """Convert PDF to Markdown (basic text extraction)."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("Install PyMuPDF: pip install pymupdf")

    doc = fitz.open(input_path)
    lines = []

    for page_num, page in enumerate(doc, 1):
        lines.append(f"## Page {page_num}")
        lines.append("")

        text = page.get_text()
        lines.append(text)
        lines.append("")
        lines.append("---")
        lines.append("")

    return "\n".join(lines)


def convert_xlsx_to_md(input_path: Path) -> str:
    """Convert Excel to Markdown tables."""
    try:
        import openpyxl
    except ImportError:
        raise ImportError("Install openpyxl: pip install openpyxl")

    wb = openpyxl.load_workbook(input_path, data_only=True)
    lines = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        lines.append(f"## {sheet_name}")
        lines.append("")

        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue

        # Filter out completely empty rows
        rows = [r for r in rows if any(c is not None for c in r)]
        if not rows:
            continue

        # Determine column count from first row
        col_count = len(rows[0])

        for i, row in enumerate(rows):
            cells = [str(c).strip() if c is not None else "" for c in row[:col_count]]
            lines.append("| " + " | ".join(cells) + " |")
            if i == 0:
                lines.append("|" + "|".join(["---"] * col_count) + "|")

        lines.append("")

    return "\n".join(lines)


def convert_file(input_path: Path, output_path: Optional[Path] = None) -> Path:
    """Convert a single file to Markdown."""
    suffix = input_path.suffix.lower()

    converters = {
        ".docx": convert_docx_to_md,
        ".doc": convert_docx_to_md,  # May not work for old .doc files
        ".pptx": convert_pptx_to_md,
        ".pdf": convert_pdf_to_md,
        ".xlsx": convert_xlsx_to_md,
        ".xls": convert_xlsx_to_md,
    }

    if suffix not in converters:
        raise ValueError(f"Unsupported file type: {suffix}. Supported: {list(converters.keys())}")

    print(f"Converting: {input_path}")
    content = converters[suffix](input_path)

    if output_path is None:
        output_path = input_path.with_suffix(".md")

    output_path.write_text(content, encoding="utf-8")
    print(f"Created: {output_path}")

    return output_path


def batch_convert(input_folder: Path, output_folder: Optional[Path] = None):
    """Convert all supported files in a folder."""
    if output_folder is None:
        output_folder = input_folder

    output_folder.mkdir(parents=True, exist_ok=True)

    extensions = [".docx", ".doc", ".pptx", ".pdf", ".xlsx", ".xls"]
    files = [f for f in input_folder.iterdir() if f.suffix.lower() in extensions]

    if not files:
        print(f"No supported files found in {input_folder}")
        return

    print(f"Found {len(files)} files to convert")

    for f in files:
        output_path = output_folder / f.with_suffix(".md").name
        try:
            convert_file(f, output_path)
        except Exception as e:
            print(f"Error converting {f}: {e}")


def main():
    parser = argparse.ArgumentParser(description="Convert documents to Markdown")
    parser.add_argument("input", help="Input file or folder (with --batch)")
    parser.add_argument("output", nargs="?", help="Output file or folder")
    parser.add_argument("--batch", action="store_true", help="Batch convert a folder")

    args = parser.parse_args()

    input_path = Path(args.input)
    output_path = Path(args.output) if args.output else None

    if not input_path.exists():
        print(f"Error: {input_path} does not exist")
        sys.exit(1)

    try:
        if args.batch:
            if not input_path.is_dir():
                print("Error: --batch requires a folder path")
                sys.exit(1)
            batch_convert(input_path, output_path)
        else:
            convert_file(input_path, output_path)
    except ImportError as e:
        print(f"Missing dependency: {e}")
        print("\nInstall required packages:")
        print("  pip install python-docx python-pptx pymupdf openpyxl")
        sys.exit(1)
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
